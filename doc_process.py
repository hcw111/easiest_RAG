"""
输入文档
读取内容，整理，分块
输出文档内容
"""
import logging
import os
import csv
import pandas as pd
from io import StringIO
from pdfminer.high_level import extract_text_to_fp
from docx import Document
from pptx import Presentation
import time

from langchain_text_splitters import RecursiveCharacterTextSplitter
logging.basicConfig(level=logging.INFO)


def extract_text(file_path):
    """负责文档读取
    :param file_path: 文档地址
    :return 如果读取成功输出文档内容，否则输出空
    记录读取日志
    """
    if not os.path.exists(file_path):
        logging.error(f"文件不存在: {file_path}")
        return
    ext = os.path.splitext(file_path)[1].lower()

    try:
        file_name = os.path.basename(file_path)
        if ext == '.txt':
            logging.info(f"正在读取txt文件: {file_name}")
            for coding in ['utf-8', 'gbk', 'ansi', 'gb2312']:
                try:
                    with open(file_path, 'r', encoding=coding) as f:
                        return f.read()
                except Exception as e:
                    continue

        elif ext == '.csv':
            logging.info(f"正在读取CSV文件: {file_name}")
            with open(file_path, 'r', newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                return '\n'.join([', '.join(row) for row in reader])

        elif ext == '.pdf':
            logging.info(f"正在读取PDF文件: {file_name}")
            output = StringIO()
            with open(file_path, 'rb') as fp:
                extract_text_to_fp(fp, output)
            return output.getvalue()

        elif ext == '.docx':
            logging.info(f"正在读取Word文件: {file_name}")
            doc = Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])

        elif ext == '.doc':
            logging.info(f"正在读取Word文件: {file_name}")
            try:
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.visible = 0
                doc = word.Documents.Open(file_path)
                text = doc.Range().Text
                doc.Close()
                word.Quit()
                return text
            except Exception as e:
                logging.error(f"[读取文件错误: {file_name}: {e}]")
                return

        elif ext in ('.xls', '.xlsx'):
            logging.info(f"正在读取Excel文件: {file_name}")
            df = pd.read_excel(file_path)
            return df.to_string()

        elif ext in ('.pptx', '.ppt'):
            logging.info(f"正在读取PPT文件: {file_name}")
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)

        else:
            logging.error(f"[不支持这种文件格式: {ext}]")
            return

    except Exception as e:
        logging.error(f"[读取文件错误: {file_path}: {e}]")
        return


def split_text(text):
    logging.info("正在分块文本...")
    try:
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=40, separators=['\n\n', '\n', '。', '，', '；', '：', ' ', ''])
        chunks = text_splitter.split_text(text)
    except Exception as e:
        logging.error(f"[分块文本错误: {e}]")
        return []

    if not chunks:
        logging.error("文档分块为空或无法对文本分块")
        return []
    return chunks


def multiple_file_process(file_paths):
    """
    处理多文件
    :param file_paths:
    :return: all_documents, all_doc_ids, all_doc_meta
    all_documents: 所有的文档内容，并分块后的结果
    all_doc_ids： 所有分块文档的id
    all_doc_meta： 所有分块文档的元数据
    """

    # 存储所有文件的文本块、id和元数据
    all_documents = []
    all_doc_ids = []
    all_doc_meta = []

    # 统计总文本块数
    total_chunks = 0

    if not file_paths:
        logging.error("上传文件为空，请上传文档")
        return [], [], []

    try:
        all_file_name = [os.path.basename(file) for file in file_paths]
        logging.info(f"接收到{len(file_paths)}个文件: {all_file_name}")

        for file_path in file_paths:
            # 利用循环对每个文档进行处理
            file_name = os.path.basename(file_path)
            # 提取文本
            text = extract_text(file_path)
            # 分块
            chunks = split_text(text)
            # 生成id和元数据
            doc_id = f"doc_{int(time.time())}_{file_name}"
            current_file_ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
            current_file_meta = [{"source": file_name, "doc_id": doc_id} for _ in chunks]

            #  逐文件保存文本块和元数据，此时每个数据按顺序一一对应。也可以用哈希表存储，但需要处理id和元数据对应问题
            all_documents.extend(chunks)
            all_doc_ids.extend(current_file_ids)
            all_doc_meta.extend(current_file_meta)

            total_chunks += len(chunks)
        logging.info(f"成功处理{len(file_paths)}个文件，共{total_chunks}个文本块")
        return all_documents, all_doc_ids, all_doc_meta

    except Exception as e:
        logging.error(f"[处理文件错误: {e}]")
        return [], [], []


if __name__ == "__main__":
    file_paths = ["C:\\Users\\hcwam\\Desktop\\不常用论文及文档\\附件2：硕士研究生学位论文开题报告和论文工作计划表-学术型 .doc", "C:\\Users\\hcwam\\Desktop\\minimind训练备注.txt"]
    chunks, ids, meta = multiple_file_process(file_paths)

